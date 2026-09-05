"""Skill: creacion de documentos de oficina (Word, Excel, PowerPoint, PDF).

Cubre el area 7 del documento de vision ("Word/Excel/PowerPoint/PDF").
Igual que cualquier otra skill, solo registra sus propias herramientas -
core/tools.py y core/agent.py no saben que esta skill existe.

Las cuatro son tier "sensitive" (escriben un archivo real en disco, igual
que write_file) - el usuario confirma la ruta exacta antes de que se cree
nada, via el mismo flujo de confirmacion que ya existe.
"""

from pathlib import Path

from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Font
from pptx import Presentation
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

from core.tools import Tool, register as register_tool, resolve_path

SKILL = {
    "name": "Creación de documentos",
    "description": "Genera archivos reales de Word, Excel, PowerPoint y PDF.",
}


def _prepared_path(raw_path: str, suffix: str) -> Path:
    path = resolve_path(raw_path)
    if path.suffix.lower() != suffix:
        path = path.with_suffix(suffix)
    path.parent.mkdir(parents=True, exist_ok=True)
    return path


async def _exec_create_word_document(arguments: dict) -> str:
    raw_path = arguments["path"]
    title = arguments.get("title", "")
    paragraphs = arguments.get("paragraphs", [])
    try:
        path = _prepared_path(raw_path, ".docx")
        doc = Document()
        if title:
            doc.add_heading(title, level=1)
        for p in paragraphs:
            if p.startswith("- "):
                doc.add_paragraph(p[2:], style="List Bullet")
            else:
                doc.add_paragraph(p)
        doc.save(str(path))
    except Exception as exc:
        return f"Error: no pude crear el documento Word: {exc}"
    return f"Documento Word creado en '{path}' ({len(paragraphs)} párrafo(s))."


async def _exec_create_excel_document(arguments: dict) -> str:
    raw_path = arguments["path"]
    sheet_name = arguments.get("sheet_name") or "Hoja1"
    headers = arguments.get("headers", [])
    rows = arguments.get("rows", [])
    try:
        path = _prepared_path(raw_path, ".xlsx")
        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name[:31]  # limite real de Excel para nombres de hoja
        if headers:
            ws.append(headers)
            for cell in ws[1]:
                cell.font = Font(bold=True)
        for row in rows:
            ws.append(row)
        wb.save(str(path))
    except Exception as exc:
        return f"Error: no pude crear el Excel: {exc}"
    return f"Excel creado en '{path}' ({len(rows)} fila(s) de datos)."


async def _exec_create_pdf_document(arguments: dict) -> str:
    raw_path = arguments["path"]
    title = arguments.get("title", "")
    paragraphs = arguments.get("paragraphs", [])
    try:
        path = _prepared_path(raw_path, ".pdf")
        styles = getSampleStyleSheet()
        story = []
        if title:
            story.append(Paragraph(title, styles["Title"]))
            story.append(Spacer(1, 14))
        for p in paragraphs:
            story.append(Paragraph(p, styles["Normal"]))
            story.append(Spacer(1, 8))
        SimpleDocTemplate(str(path), pagesize=LETTER).build(story)
    except Exception as exc:
        return f"Error: no pude crear el PDF: {exc}"
    return f"PDF creado en '{path}' ({len(paragraphs)} párrafo(s))."


async def _exec_create_powerpoint_document(arguments: dict) -> str:
    raw_path = arguments["path"]
    title = arguments.get("title", "")
    slides = arguments.get("slides", [])
    try:
        path = _prepared_path(raw_path, ".pptx")
        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title or "Presentación"

        content_layout = prs.slide_layouts[1]
        for slide_spec in slides:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = slide_spec.get("heading", "")
            body = slide.placeholders[1].text_frame
            bullets = slide_spec.get("bullets", [])
            body.clear()
            for i, bullet in enumerate(bullets):
                paragraph = body.paragraphs[0] if i == 0 else body.add_paragraph()
                paragraph.text = bullet
        prs.save(str(path))
    except Exception as exc:
        return f"Error: no pude crear el PowerPoint: {exc}"
    return f"PowerPoint creado en '{path}' ({len(slides)} diapositiva(s) de contenido)."


def register() -> None:
    register_tool(Tool(
        name="create_word_document",
        description="Crea un documento de Word (.docx) real con un título y una lista de párrafos.",
        parameters={
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Ruta del archivo .docx a crear."},
                "title": {"type": "string", "description": "Título del documento (opcional)."},
                "paragraphs": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Párrafos del documento en orden. Un párrafo que empiece con '- ' se agrega como viñeta.",
                },
            },
            "required": ["path", "paragraphs"],
        },
        tier="sensitive",
        executor=_exec_create_word_document,
        confirm_text=lambda a: f"Crear documento Word: {a.get('path')}",
    ))

    register_tool(Tool(
        name="create_excel_document",
        description="Crea una hoja de cálculo de Excel (.xlsx) real con encabezados y filas de datos.",
        parameters={
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Ruta del archivo .xlsx a crear."},
                "sheet_name": {"type": "string", "description": "Nombre de la hoja (opcional)."},
                "headers": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Encabezados de columna (opcional).",
                },
                "rows": {
                    "type": "array",
                    "items": {"type": "array", "items": {}},
                    "description": "Filas de datos, cada una es una lista de valores de celda.",
                },
            },
            "required": ["path", "rows"],
        },
        tier="sensitive",
        executor=_exec_create_excel_document,
        confirm_text=lambda a: f"Crear Excel: {a.get('path')}",
    ))

    register_tool(Tool(
        name="create_pdf_document",
        description="Crea un documento PDF real con un título y una lista de párrafos.",
        parameters={
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Ruta del archivo .pdf a crear."},
                "title": {"type": "string", "description": "Título del documento (opcional)."},
                "paragraphs": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Párrafos del documento en orden.",
                },
            },
            "required": ["path", "paragraphs"],
        },
        tier="sensitive",
        executor=_exec_create_pdf_document,
        confirm_text=lambda a: f"Crear PDF: {a.get('path')}",
    ))

    register_tool(Tool(
        name="create_powerpoint_document",
        description="Crea una presentación de PowerPoint (.pptx) real con una diapositiva de título y diapositivas de contenido con viñetas.",
        parameters={
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Ruta del archivo .pptx a crear."},
                "title": {"type": "string", "description": "Título de la diapositiva inicial."},
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "heading": {"type": "string"},
                            "bullets": {"type": "array", "items": {"type": "string"}},
                        },
                    },
                    "description": "Diapositivas de contenido, cada una con un título y una lista de viñetas.",
                },
            },
            "required": ["path", "slides"],
        },
        tier="sensitive",
        executor=_exec_create_powerpoint_document,
        confirm_text=lambda a: f"Crear PowerPoint: {a.get('path')}",
    ))
